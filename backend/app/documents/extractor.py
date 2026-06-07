from pathlib import Path
from typing import Protocol

import fitz
import pdfplumber
from docx import Document
from pptx import Presentation


class TextExtractor(Protocol):
    def supports(self, path: Path) -> bool:
        """Return true when the extractor can process this file."""

    async def extract_text(self, path: Path) -> str:
        """Extract readable study text from a file."""


class UnsupportedDocumentError(ValueError):
    pass


class PlainTextExtractor:
    _extensions = {".txt", ".md"}

    def supports(self, path: Path) -> bool:
        return path.suffix.lower() in self._extensions

    async def extract_text(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore").strip()


class PdfExtractor:
    def supports(self, path: Path) -> bool:
        return path.suffix.lower() == ".pdf"

    async def extract_text(self, path: Path) -> str:
        text = self._extract_with_pymupdf(path)
        if text.strip():
            return text.strip()
        return self._extract_with_pdfplumber(path).strip()

    def _extract_with_pymupdf(self, path: Path) -> str:
        chunks: list[str] = []
        with fitz.open(path) as document:
            for page in document:
                chunks.append(page.get_text())
        return "\n".join(chunks)

    def _extract_with_pdfplumber(self, path: Path) -> str:
        chunks: list[str] = []
        with pdfplumber.open(path) as document:
            for page in document.pages:
                chunks.append(page.extract_text() or "")
        return "\n".join(chunks)


class DocxExtractor:
    def supports(self, path: Path) -> bool:
        return path.suffix.lower() == ".docx"

    async def extract_text(self, path: Path) -> str:
        document = Document(path)
        return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()


class PptxExtractor:
    def supports(self, path: Path) -> bool:
        return path.suffix.lower() == ".pptx"

    async def extract_text(self, path: Path) -> str:
        presentation = Presentation(path)
        chunks: list[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            slide_chunks = [f"Diapositiva {index}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_chunks.append(shape.text)
            chunks.append("\n".join(slide_chunks))

        return "\n\n".join(chunks).strip()


class DocumentTextExtractor:
    def __init__(self, extractors: list[TextExtractor] | None = None) -> None:
        self._extractors = extractors or [
            PlainTextExtractor(),
            PdfExtractor(),
            DocxExtractor(),
            PptxExtractor(),
        ]

    async def extract_text(self, path: Path) -> str:
        for extractor in self._extractors:
            if extractor.supports(path):
                return await extractor.extract_text(path)

        supported = ", ".join([".txt", ".md", ".pdf", ".docx", ".pptx"])
        raise UnsupportedDocumentError(
            f"Unsupported document type '{path.suffix}'. Supported types: {supported}."
        )

    def supports(self, path: Path) -> bool:
        return any(extractor.supports(path) for extractor in self._extractors)
